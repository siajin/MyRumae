"""Text-only extraction for Office formats (.pptx, .docx).

Per the Phase 1 decision, OCR / image / shape geometry are out of scope —
we only need plain text so the DOCX summary's "파싱 본문" section has
something to show for slide decks and Word handouts. Output uses the same
ParseResult / Block shape as parser/pdf.py, so downstream code (JSON
write, ParsedContent upsert) is unchanged.

Slide-level page numbering for pptx (1-based) matches what users see in
PowerPoint. docx has no page concept here, so we keep page=1 for every
paragraph.
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

from .blocks import Block, ParseResult

log = logging.getLogger(__name__)


def _pptx_version() -> str:
    try:
        import pptx
        return getattr(pptx, "__version__", "unknown")
    except Exception:
        return "unavailable"


def _docx_version() -> str:
    try:
        import docx
        return getattr(docx, "__version__", "unknown")
    except Exception:
        return "unavailable"


def parser_version(ext: str) -> str:
    if ext == "pptx":
        return f"python-pptx-{_pptx_version()}"
    if ext == "docx":
        return f"python-docx-{_docx_version()}"
    return f"office-{ext}"


def _extract_pptx(path: Path) -> ParseResult:
    from pptx import Presentation

    pres = Presentation(str(path))
    blocks: list[Block] = []
    slide_count = 0

    for slide_idx, slide in enumerate(pres.slides, start=1):
        slide_count = slide_idx
        order = 0
        for shape in slide.shapes:
            text = ""
            if shape.has_text_frame:
                # Join paragraph runs with newline so list/title structure
                # survives in plain text.
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if not line:
                        line = (para.text or "")
                    if line.strip():
                        paragraphs.append(line)
                text = "\n".join(paragraphs).strip()
            elif getattr(shape, "has_table", False) and shape.has_table:
                # Flatten table cells row-by-row.
                rows = []
                for row in shape.table.rows:
                    cells = [(c.text or "").strip() for c in row.cells]
                    rows.append("\t".join(cells))
                text = "\n".join(r for r in rows if r.strip()).strip()

            if not text:
                continue

            order += 1
            blocks.append(
                Block(page=slide_idx, kind="text", text=text, order=order)
            )

    return ParseResult(
        blocks=blocks,
        page_count=slide_count,
        used_ocr=False,
        parser_version=parser_version("pptx"),
    )


def _extract_docx(path: Path) -> ParseResult:
    from docx import Document

    doc = Document(str(path))
    blocks: list[Block] = []
    order = 0

    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        order += 1
        blocks.append(Block(page=1, kind="text", text=text, order=order))

    # Tables in body — flatten cells.
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            rows.append("\t".join(cells))
        text = "\n".join(r for r in rows if r.strip()).strip()
        if not text:
            continue
        order += 1
        blocks.append(Block(page=1, kind="table", text=text, order=order))

    return ParseResult(
        blocks=blocks,
        page_count=1 if blocks else 0,
        used_ocr=False,
        parser_version=parser_version("docx"),
    )


def extract_office(
    path: Path,
    *,
    material_id: Optional[int] = None,
) -> ParseResult:
    """Dispatch by extension. Caller (pipeline.parse_material) handles
    failures via the surrounding try/except."""
    ext = path.suffix.lstrip(".").lower()
    if ext == "pptx":
        return _extract_pptx(path)
    if ext == "docx":
        return _extract_docx(path)
    raise ValueError(f"unsupported office ext: {ext}")
